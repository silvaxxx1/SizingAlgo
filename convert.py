import os
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
from weasyprint import HTML
from pdf2image import convert_from_path

# --- CONFIG ---
INPUT_HTML = "v2g_presentation.html"  # your HTML file
OUTPUT_PPTX = "slides.pptx"
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
TEMP_DIR = "temp_slides"

# --- PREP ---
os.makedirs(TEMP_DIR, exist_ok=True)
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# --- PARSE HTML ---
with open(INPUT_HTML, "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f, "html.parser")

# --- FIND SLIDES ---
slides = soup.find_all(["section", "div"], class_="slide") or soup.find_all("section")
if not slides:
    slides = [soup]  # fallback: treat entire HTML as one slide

print(f"Found {len(slides)} slides.")

# --- RENDER EACH SLIDE ---
for i, slide_html in enumerate(slides, 1):
    slide_file = os.path.join(TEMP_DIR, f"slide_{i}.html")
    pdf_file = os.path.join(TEMP_DIR, f"slide_{i}.pdf")
    img_file = os.path.join(TEMP_DIR, f"slide_{i}.png")

    # Save individual HTML fragment
    with open(slide_file, "w", encoding="utf-8") as f:
        f.write(str(slide_html))

    # Render HTML → PDF
    HTML(slide_file).write_pdf(pdf_file)

    # Convert PDF → PNG (first page only)
    images = convert_from_path(pdf_file, dpi=200)
    images[0].save(img_file, "PNG")

    # Add slide to PPTX
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(img_file, 0, 0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT)
    print(f"Rendered slide {i}")

# --- SAVE PRESENTATION ---
prs.save(OUTPUT_PPTX)
print(f"✅ Presentation saved as {OUTPUT_PPTX}")
